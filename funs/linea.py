from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt


def lineas_ap_glob(prs, nro_slide, apr_glob):
    # Definir slide
    slide = prs.slides[nro_slide-1]

    for i in [(8.46, 15.64), (25.5, 32.44)]:

        # Mínimo y máximo del cuadro del gráfico
        min_x_org, max_x_org = i
        rango = max_x_org - min_x_org

        # Ubicación de la linea
        ap_glob = min_x_org+(rango*apr_glob)

        # Defición de las coordenadas
        x, y = ap_glob, 6.5

        largo = 10.8   # Largo estandar de la linea
        ancho = 2      # Ancho estandar de la linea

        line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Cm(x), Cm(y), Cm(0), Cm(largo))
        line.line.color.rgb = RGBColor(55, 96, 146)
        line.line.width = Pt(ancho)