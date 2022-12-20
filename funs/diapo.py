import pandas as pd

from funs.dimensiones import *
from funs.porcent import porcent, per_str
from funs.items import items
from funs.areas import dict_areas
from funs.tabla import *
from funs.car import *
from funs.otros import get_unique

from data.texto.dimensiones import dict_dims

from pptx.chart.data import ChartData, CategoryChartData


def check_diapo(prs, nro_slide=int, print_text=False):
    """
    Imprime el tipo de elementos que existe en una diapositiva

    prs   = Presentación
    nro_slide = Número de la diapositiva
    """
    count = 0
    slide = prs.slides[nro_slide-1]
    for shape in slide.shapes:
        if print_text == True:
            if shape.has_text_frame:
                print(count, shape.text)
            else:
                print(count, shape)
        else:
            print(count, shape)

        count += 1

def diapo_replace_txt(prs, data, nro_slide=int, txts=list):
    """
    Reemplaza los textos, por los valores asociados a cada área

    prs = Template
    data = Diccionario de los datos
    nro_slide = Número de la diapositiva
    txts = Lista de el/los texto(s) que se desean reemplazar en el template

    * El nombre el texto a reemplazar tiene que ser el mismo que el que se encuentra en el diccionario donde está la data
    """
    slide = prs.slides[nro_slide-1]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for txt in txts:
            text_frame.text = text_frame.text.replace(txt, data[txt])

def del_diapo(prs, area=str, areas_glob=list, sub_areas=list, unico=bool):
    """
    Elimina las diapositivas específicas
    Puede ser una, o un rango determinado
    Se puede definir la cantidad de subareas mínimas

    prs = Presentación
    sub_areas = Lista de subareas del segmento continuo
    min_areas = Cantidad mínima de áreas para eliminar ciertas diapos
    """
    
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)

    if area != 'Global Vitacura':
        xml_slides.remove(slides[68-1])

    if area not in areas_glob:
        # Resultados por dimensiones específicas
        xml_slides.remove(slides[67-1])

    if area != 'Global Vitacura':
        # Oportunidades de mejora desde lo inferencial
        xml_slides.remove(slides[65])
    
    if area not in areas_glob:    
        # Pregunta abierta y Resultados inferenciales
        for r in range(50, 62):
            xml_slides.remove(slides[r])
    
    if area == 'Área Educación':
        # Nube de palabras
        for r in range(57, 60):
            xml_slides.remove(slides[r+1])
        # Inferenciales
        for r in range(51, 54):
            xml_slides.remove(slides[r+1])

    if area == 'Área Salud':
        # Nube de palabras
        xml_slides.remove(slides[61])
        for r in range(57, 59):
            xml_slides.remove(slides[r+1])
        # Inferenciales
        xml_slides.remove(slides[55])
        for r in range(51, 53):
            xml_slides.remove(slides[r+1])

    if area == 'Área Municipal':
        # Nube de palabras
        for r in range(59, 61):
            xml_slides.remove(slides[r+1])
        xml_slides.remove(slides[58])
        
        # Inferenciales
        for r in range(53, 55):
            xml_slides.remove(slides[r+1])
        xml_slides.remove(slides[53])

    # Tipo de jornada no va en ninguna
    if area in areas_glob:
        xml_slides.remove(slides[38+1])
        xml_slides.remove(slides[39+1])

    # Diapos de características
    if area not in areas_glob:
        for r in range(31, 43):
            xml_slides.remove(slides[r])
        if len(sub_areas) == 0:
            for r in range(27, 31):
                xml_slides.remove(slides[r])

    # Comparación de ent idc y sat con años anteriores
    if unico == True:
        xml_slides.remove(slides[24])
    
    # Comparación con otros estudios de clima
    if area != 'Global Vitacura':
        xml_slides.remove(slides[21])

    # Comparación con años anteriores
    if unico == True:
        xml_slides.remove(slides[17])
        
    # Comparación segun sobarea
    if area == 'Global Vitacura':
        xml_slides.remove(slides[16])
    
    # Indices
    if area in areas_glob:
        xml_slides.remove(slides[4-1])
        xml_slides.remove(slides[3-1])
    else:
        if len(sub_areas) > 0:
            xml_slides.remove(slides[3-1])
        else:
            xml_slides.remove(slides[4-1])

        xml_slides.remove(slides[2-1])


    



def diapo_dotacion(prs, data, nro_slide=int, nro_shape=int):
    """
    Función específica para mostrar la tasa de respuesta
    Asume que '<c>' y '<nc>' se encuentran como key en el diccionario.

    prs = Template
    data = Diccionario de los datos
    nro_slide = Número de la diapositiva
    nro_shape = Número del elemento que corresponde al gráfico

    """
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = ChartData()
    chart_data.categories = ['Encuestas Contestadas', 'Encuestas No Contestadas']
    c =  porcent(data['<asi_yes>'], data['<asi_not>'])
    nc = porcent(data['<asi_not>'], data['<asi_yes>'])
    chart_data.add_series('Dotación', (c, nc), number_format='0%')
    chart.replace_data(chart_data)

def diapo_pie(prs, data, nro_slide=int, nro_shape=int, key=str):
    """
    Reemplaza los datos para hacer una pie de la distribución de los valores positivos neutro y negativos de manera global

    prs = Template
    data = Diccionario de los datos
    nro_slide = Número de la diapositiva
    nro_shape = Número del elemento que corresponde al gráfico

    """
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = ChartData()
    chart_data.categories = ['Apreciación Positiva', 'Apreciación Neutra', 'Apreciación Negativa']
    # if glob == True:
    #     chart_data.add_series('Apreciación', (data['<glob_pos>'], data['<glob_neu>'], data['<glob_neg>']), number_format='0%')
    # if sob_area == True:
    #     chart_data.add_series('Apreciación', (data['<sob_area_pos>'], data['<sob_area_neu>'], data['<sob_area_neg>']), number_format='0%')
    # else:
    #     chart_data.add_series('Apreciación', (data['<area_pos>'], data['<area_neu>'], data['<area_neg>']), number_format='0%')
    chart_data.add_series('Apreciación', (data[f"<{key}_pos>"], data[f"<{key}_neu>"], data[f"<{key}_neg>"]), number_format='0%')
    
    chart.replace_data(chart_data)

def diapo_dims(prs, data, nro_slide=int, nro_shape=int, dims=list, key=str, signo=False, comp=False, str_comp=str, comp2=False, str_comp2=str):
    """
    Reemplaza los datos del gráfico seleccionado por los datos del diccionario.
    Se puede elegir si se quieren los restultados positivos, neutro, negativos o todos los anteriores

    prs = Template
    data = Diccionario de los datos
    nro_slide = Número de la diapositiva
    nro_shape = Número del elemento que corresponde al gráfico
    dims = Lista de dimensiones elegidas
    key = Determina qué dimensión o dimensiones se reemplazarán, para esto, consulta a los datos del diccionario

    """
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()
    chart_data.categories = dims
    if signo == False:
        chart_data.add_series('Apreciación Positiva', data[f'<{key}_pos>'], number_format='0%')
        chart_data.add_series('Apreciación Neutra',   data[f'<{key}_neu>'], number_format='0%')
        chart_data.add_series('Apreciación Negativa', data[f'<{key}_neg>'], number_format='0%')
    elif type(signo) == str:
        chart_data.add_series('Apreciación Positiva', data[f'<{key}_{signo}>'], number_format='0%')
        
        if type(comp) != bool:
            if type(comp) == list:
                chart_data.add_series(str_comp, comp, number_format='0%')
            else:
                chart_data.add_series('Apreciación Positiva Global', [comp for i in dims], number_format='0%')
        
        if type(comp2) != bool:
            chart_data.add_series(str_comp2, comp2, number_format='0%')


    chart.replace_data(chart_data)


def diapo_items(prs, data, nro_slide=int, nro_shape=int, dict_items=dict, item=str):
    """
    Reemplaza los datos del gráfico por la distribución de puntajes de cada ítem de la dimensión elegida.

    prs = Template
    data = Diccionario de los datos
    dict_items = Diccionario donde se encuentran el código de los ítems con los respectivos enunciados
    nro_slide = Número de la diapositiva
    nro_shape = Número del elemento que corresponde al gráfico
    item = La sigla del la dimensión que se quiere analizar por items. Este tiene que estar definido de manera previa en el diccionario

    """
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()
    chart_data.categories = items(dict=dict_items, dim=item)
    chart_data.add_series('Apreciación Positiva', data[f'<{item}_item_pos>'], number_format='0%')
    chart_data.add_series('Apreciación Neutra',   data[f'<{item}_item_neu>'], number_format='0%')
    chart_data.add_series('Apreciación Negativa', data[f'<{item}_item_neg>'], number_format='0%')
    chart.replace_data(chart_data)

def diapo_area(prs, df, dims, nro_slide=int, nro_shape=int, nxt_seg=str, sub_areas=list, areas_not=list):
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()
    sub_areas = [i for i in sub_areas if i not in areas_not]
    if len(sub_areas) != 0:
        chart_data.categories = sub_areas
        chart_data.add_series('Apreciación Positiva', dict_areas(df, nxt_seg, sub_areas, dims)['pos'], number_format='0%')
        chart_data.add_series('Apreciación Positiva', dict_areas(df, nxt_seg, sub_areas, dims)['neu'], number_format='0%')
        chart_data.add_series('Apreciación Positiva', dict_areas(df, nxt_seg, sub_areas, dims)['neg'], number_format='0%')
        chart.replace_data(chart_data)

def diapo_car(prs, df, nro_car=int, dims=list, nro_slide=int, nro_shape=int, orden=False):
    """
    Rellena los datos de los gráficos considerando las características de los participantes
    Cada grupo de características está separado y agrupado de manera diferente.
    
    prs = presentación
    df = base de datos
    nro_car = número de la característica a analizar
    dims = Dimensiones considerarán
    nro_slide = Número de la diapositiva
    nro_shape = Número del objeto (gráfico) que se le van a integrar los datos
    orden = Lista con las características ordenadas de manera previa
    """

    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()

    if orden == False:
        categorias = get_unique(df, column=f'car_{nro_car}')
        # categorias = [i for i in pd.unique(df[f'car_{nro_car}']) if type(i) == str]
    else:
        categorias = [i for i in orden if i in pd.unique(df[f'car_{nro_car}'])]

    chart_data.categories = categorias
    dict_cars = dict_car(df, nro_car=nro_car, dims=dims, list_sort=categorias)
    chart_data.add_series('Apreciación Positiva', dict_cars['pos'], number_format='0%')
    chart_data.add_series('Apreciación Neutra',   dict_cars['neu'], number_format='0%')
    chart_data.add_series('Apreciación Negativa', dict_cars['neg'], number_format='0%')
    chart.replace_data(chart_data)


def diapo_tabla(prs, df, nro_slide=int, nro_shape=int, dims=list, titulo=str, nro_car=False, nxt_seg=False, sub_areas=False, areas_not=list, heath_map=True, orden=False):
    """
    Toma elementos y los agrega en una tabla preexistente.
    Esa tabal tiene solo 3 filas que irán aumentando según la cantidad de agrupaciones
    La función aguanta tanto características como subareas
    """
    slide = prs.slides[nro_slide-1]
    table = slide.shapes[nro_shape].table

    # Si la tabla generada va a ser de una característica puntual
    if type(nro_car) == int:
        if orden == False:
            idx = get_unique(df, column=f"car_{nro_car}")
        elif type(orden) == list:
            idx = orden
    
    # Si la tabla generada va a ser de las subareas
    elif nxt_seg and sub_areas != False:
        sub_areas = [i for i in sub_areas if i not in areas_not] # Este filtro no debería estar, debería estar por defecto más adelante
        idx = [i for i in sub_areas if len(df[df[nxt_seg] == i]) >= 7] # Faltaría aregra el criterio (tasa_respuesta >= 0.5)
    

    # Cantidad de filas según la cantidad de elementos (característica o subareas)
    len_idx = len(idx)

    for r in range(len_idx-1):
        add_row(table)

    # Nombres de la celda del grupo
    cell = table.cell(1, 0)
    cell_set(cell, cell_text=titulo, text_color=RGBColor(0, 0, 0), bold=True)

    # Ingresar valores de la 2da fila
    for d in range(len(dims)):

        # Segunda Fila (los promedios de cada dimensión)
        cell = table.cell(1, d+1)
        glob = get_per_dims(df, dims=[dims[d]], signo='pos', promedio=True, string=True)
        cell_set(cell, cell_text=glob, text_color=RGBColor(0, 0, 0), bold=True)
    
    for i in range(len_idx):
        # Ingresa el nombre de las características/subareas en la primera columna
        cell = table.cell(i+2, 0)
        cell_set(cell, cell_text=idx[i], text_color=RGBColor(0, 0, 0))
        
        for j in range(len(dims)):
            # Ingresar el valor de las diferentes dimensiones
            cell = table.cell(i+2, j+1)
            
            # La columna que se elija dependerá de si se está comparando una característica, o una subarea
            if type(nro_car) == int:
                columna = f'car_{nro_car}'
            elif nxt_seg and sub_areas != False:
                columna = nxt_seg

            val = get_per_dims(df[df[columna] == idx[i]], dims=[dims[j]], signo='pos', promedio=True)
            val_str = per_str(val)
            cell_set(cell, cell_text=val_str, text_color=RGBColor(0, 0, 0))

            # Teñir celdas según sus valores
            if heath_map == True:
                por_glo = get_per_dims(df, dims=[dims[j]], signo='pos', promedio=True)
                por_raw = val
                heath(cell, glo=por_glo, raw=por_raw)


def diapo_rank(prs, nro_slide=int, nro_shape=int, tops=dict, dict_items=dict):
    
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()

    categorias = [f"[{t[:3].upper()}] {dict_items[t]}" for t in tops.keys()]
    categorias.reverse()

    chart_data.categories = categorias
    values = [t for t in tops.values()]
    values.reverse()

    chart_data.add_series('Tops', values, number_format='0%')
    chart.replace_data(chart_data)

def diapo_ir(prs, nro_slide=int, nro_shape=int, dict_dims=dict):
    
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()

    categorias = [dict_dims[k] for k in list(dict_dims.keys())]
    categorias.reverse()
    chart_data.categories = categorias
    
    values = [v for v in list(dict_dims.values())]
    values.reverse()

    chart_data.add_series('Importancia Relativa', values, number_format='0%')
    chart.replace_data(chart_data)

def diapo_ir_txt(prs, nro_slide=int, nro_shape=int, all_dims=dict, tops=dict, glob_mean=list):
    """
    Agregar las dimensiones más elegidas según la importancia relativa (ranking) además de agregar su porcentaje
    La celda del porcentaje se tiñe si es que la apreciación positiva de la persona se encuentra sobre el promedio o bajo el promedio (para entender mejor, ver L16)

    all_dims = Diccionario donde se encuentran todas las dimenciones, con su porcentaje de apreciación positiva. De este diccionario se obtendrán los valores que se compararán con el 'glob_mean'
    tops = Diccionario con las dimensiones con mayor "importancia relativa" en orden de mayor a menor
    glob_mean = Promedio de aprecición positiva del área. Esto servirá para teñir las celdas en caso de que las dimensiones del 'tops' tengan una apreciación positiva menor al promedio de su área.
    """

    dims_tops = [dict_dims[i] for i in list(tops.keys())]
    per_str_ir = [per_str(i) for i in tops.values()]
    heath_list = [all_dims[i] for i in tops.keys()]

    slide = prs.slides[nro_slide-1]
    table = slide.shapes[nro_shape].table

    for i in range(len(dims_tops)):
        cell = table.cell(i+1, 1)
        cell_set(cell, cell_text=dims_tops[i], alignment_center=False, text_size=22)

    for i in range(len(per_str_ir)):
        cell = table.cell(i+1, 2)
        # En caso de ser mayor a la puntuación global
        if heath_list[i] > glob_mean:
            cell_set(cell, cell_text='+', alignment_center=True, text_size=24, text_color=RGBColor(0, 154, 70), bold=True)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(215, 228, 189)
        # En caso de ser menor a la puntuación global
        else:
            cell_set(cell, cell_text='-', alignment_center=True, text_size=24, text_color=RGBColor(192, 0, 0), bold=True)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(253, 212, 181)


def diapo_dims_ent_com_sat(prs, data1, data2, nro_slide=int, nro_shape=int, dims=list, signo=False):
    """
    Reemplaza los datos del gráfico seleccionado por los datos del diccionario.
    Se puede elegir si se quieren los restultados positivos, neutro, negativos o todos los anteriores

    prs = Template
    data = Diccionario de los datos
    nro_slide = Número de la diapositiva
    nro_shape = Número del elemento que corresponde al gráfico
    dims = Lista de dimensiones elegidas
    key = Determina qué dimensión o dimensiones se reemplazarán, para esto, consulta a los datos del diccionario

    """
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()
    chart_data.categories = dims

    valores_2021 = [data1[f'<enl_{signo}>']] + [data1[f'<idc_{signo}>']] + [data1[f'<sat_{signo}>']]
    valores_2022 = [data2[f'<enl_{signo}>']] + [data2[f'<idc_{signo}>']] + [data2[f'<sat_{signo}>']]

    chart_data.add_series('2021', valores_2021, number_format='0%')
    chart_data.add_series('2022', valores_2022, number_format='0%')

    chart.replace_data(chart_data)


def diapo_comp_bars(prs, data1, data2, nro_slide=int, nro_shape=int, key=list, series=list):
    """
    Ingresa los datos de 2 series de barras para que se comparen entre sí.
    Permite comparar resultados según año, o según sub/sobgrupo.

    prs = Template
    data = Diccionario de los datos
    nro_slide = Número de la diapositiva
    nro_shape = Número del elemento que corresponde al gráfico
    key = Lista de las llaves del diccionario que se desean comparar (tiene que tener 2 valores siempre)
    series = Lista de nombres que se les dará a cada una de las 2 series
    """
    slide = prs.slides[nro_slide-1]
    chart = slide.shapes[nro_shape].chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Apreciación Positiva', 'Apreciación Neutra', 'Apreciación Negativa']

    valores_1 = [data1[f'<{key[0]}_pos>']] + [data1[f'<{key[0]}_neu>']] + [data1[f'<{key[0]}_neg>']]
    valores_2 = [data2[f'<{key[1]}_pos>']] + [data2[f'<{key[1]}_neu>']] + [data2[f'<{key[1]}_neg>']]

    chart_data.add_series(series[0], valores_1, number_format='0%')
    chart_data.add_series(series[1], valores_2, number_format='0%')

    chart.replace_data(chart_data)


def diapo_cor(prs, nro_slide, dict_cor):
    slide = prs.slides[nro_slide-1]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        
        if text_frame.text[0] == "<":
            text_frame.text = text_frame.text.replace(text_frame.text, str(dict_cor[text_frame.text]))

    prs.save("ppts/autorreportes/test.pptx")